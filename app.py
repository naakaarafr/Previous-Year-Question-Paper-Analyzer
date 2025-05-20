import streamlit as st
import io
import os
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import tempfile
from pptx import Presentation   # For PPT files
import pytesseract  # For OCR

# Load environment variables and configure API
load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
if api_key:
    genai.configure(api_key=api_key)
else:
    st.error("Google API Key not found. Please set the GOOGLE_API_KEY environment variable.")

# Check for PDF2Image availability
try:
    from pdf2image import convert_from_bytes, convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

def extract_text_from_pdf(pdf_bytes):
    """Extract text from PDF using PyPDF2 and OCR if needed"""
    text = ""
    
    # Try PyPDF2 first
    try:
        pdf_reader = PdfReader(io.BytesIO(pdf_bytes))
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            text += page_text + "\n" if page_text else "\n"
    except Exception as e:
        st.warning(f"PyPDF2 extraction issue: {e}. Trying alternative method...")
    
    # If text is empty or insufficient, try OCR-based extraction
    if not text.strip() or len(text.strip()) < 100:
        if PDF2IMAGE_AVAILABLE:
            try:
                with st.spinner("Using OCR to extract text from images in PDF..."):
                    text = extract_text_with_ocr(pdf_bytes)
            except Exception as e:
                st.error(f"OCR extraction failed: {e}")
        else:
            st.error("PDF contains images but pdf2image is not available. Install pdf2image and poppler for OCR.")
    
    return text

def extract_text_with_ocr(pdf_bytes):
    """Extract text from PDF using OCR"""
    if not PDF2IMAGE_AVAILABLE:
        raise ImportError("pdf2image and poppler are required for OCR")
    
    try:
        # Convert PDF to images
        images = convert_from_bytes(
            pdf_bytes,
            dpi=300,  # Higher DPI for better OCR results
            fmt='jpeg',
            thread_count=1,
            strict=False,
            use_cropbox=True,
            transparent=False
        )
        
        if not images or len(images) == 0:
            raise Exception("No images extracted from PDF. The document may be empty or corrupted.")
        
        # Extract text from each image using OCR
        extracted_content = ""
        
        for i, image in enumerate(images):
            # Add page marker
            extracted_content += f"\n\n--- Page {i + 1} ---\n\n"
            
            # Perform OCR on the image
            try:
                # Use pytesseract for OCR
                image_text = pytesseract.image_to_string(image)
                
                # If OCR text is empty, note it
                if not image_text.strip():
                    extracted_content += f"[Image content on page {i + 1} - no text detected]\n"
                else:
                    extracted_content += image_text + "\n"
            except Exception as ocr_error:
                st.warning(f"OCR error on page {i+1}: {str(ocr_error)}")
                extracted_content += f"[OCR failed for page {i + 1}]\n"
        
        return extracted_content
    except Exception as e:
        raise Exception(f"Error extracting with OCR: {str(e)}")

def extract_text_from_ppt(ppt_bytes):
    """Extract text from PowerPoint files"""
    text = ""
    # Save the bytes to a temporary file
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
        temp_file.write(ppt_bytes)
        temp_file_path = temp_file.name
    
    try:
        # Open the presentation
        presentation = Presentation(temp_file_path)
        
        # Extract text from all slides
        for i, slide in enumerate(presentation.slides):
            text += f"\n--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        st.error(f"Error extracting text from PPT: {e}")
    finally:
        # Clean up the temporary file
        if os.path.exists(temp_file_path):
            os.unlink(temp_file_path)
    
    return text

def get_text_chunks(text):
    """Split text into chunks for processing"""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=10000,
        chunk_overlap=1000
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    """Create and save vector store from text chunks"""
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")
    return vector_store

def get_conversational_chain():
    """Create a QA chain for answering questions"""
    prompt_template = """
    Role & Purpose:
    You are an AI-powered Previous Year Question Paper Analyzer designed to assist students, educators, and researchers in understanding trends, patterns, and insights from historical exam papers. Your core function is to process multiple PDF question papers, analyze their content, answer questions directly from the papers, and provide actionable, data-driven responses. Prioritize accuracy, clarity, and educational value in all interactions.

    Core Capabilities:

    Multi-PDF Processing & Analysis

    Accept and analyze multiple PDFs (e.g., 5+ years of question papers) to identify cross-year patterns.

    Extract metadata (subject, year, board/university, marks distribution).

    Parse text, tables, diagrams, and mathematical content (if OCR-supported).

    Key Analytical Tasks

    Topic Frequency Analysis: Identify recurring topics/subtopics, their weightage, and frequency across years.

    Difficulty Trend Mapping: Classify questions by difficulty (easy/medium/hard) and track changes over time.

    Question-Type Categorization: Highlight common formats (MCQs, essays, numericals) and their distribution.

    Mark Distribution: Break down marks allocation per chapter, unit, or topic.

    Pattern Prediction: Suggest probable questions/chapters based on historical trends (with disclaimers).

    Question Answering

    Direct Question Resolution: Accurately answer questions from the uploaded papers by:

    Extracting the exact question text (verbatim) from the PDFs.

    Providing step-by-step solutions for numerical/theoretical questions.

    Highlighting correct answers for MCQs with explanations.

    Offering model essay outlines or key points for descriptive questions.

    Contextual Support: If diagrams, tables, or equations are missing due to OCR limitations, flag them and request clarification.

    Source Citation: Always reference the source paper (year, question number, page) when answering.

    User Interaction Workflow

    Step 1: Confirm receipt of uploaded PDFs and initiate analysis.

    Step 2: Generate a concise summary report (optional) unless the user opts out.

    Example: *"Analyzed 10 CBSE Class 12 Maths papers (2014–2024). Key trends: Calculus (35% weightage), 6+ mark differential equations yearly."*

    Step 3: Invite user queries with examples:

    "Which topics have the highest frequency in the last 5 years?"

    "Show difficulty trends for Organic Chemistry questions."

    "What is the answer to Question 7(a) from the 2020 paper?"

    "Explain the solution to the MCQ on page 3 of the 2022 document."

    Response Guidelines:

    Accuracy: Cite specific papers/years (e.g., "2022 Paper, Q7") to support answers.

    Visualization: Use tables, charts, or bullet points to simplify complex trends (e.g., mark distribution matrices).

    Step-by-Step Solutions: Break down answers for numerical/theoretical questions into logical steps (e.g., "Step 1: Apply Ohm's Law…").

    Limitations: Acknowledge gaps (e.g., "Diagram missing in scanned PDF; unable to answer part (b)").

    Educational Tone: Explain concepts if needed (e.g., "Differential equations often require integrating multiple steps, which explains their high marks").

    Ambiguity Handling: Clarify vague queries (e.g., "Did you mean Question 5 from the 2019 paper or 2021 paper?").

    Privacy & Compliance:

    Do not store, share, or retain uploaded PDFs post-analysis.

    Comply with data protection regulations (e.g., FERPA, GDPR) for educational content.

    Closing Interaction:

    End with a prompt for further assistance (e.g., "What else would you like to explore?").

    Politely decline non-relevant requests (e.g., "I specialize in paper analysis and question resolution only").

    Example Interaction:
    User: "What is the answer to Question 3(b) in the 2023 Physics paper?"
    AI: *"From the 2023 Physics Paper (Page 5, Question 3b):
    Question: 'Calculate the electric potential at a point 2m from a 5μC charge.'
    Solution:
    1. Use the formula for electric potential:
    V = kQ/r
    where k = 9 × 10^9 N m²/C² (Coulomb's constant), Q = 5μC = 5 × 10^−6 C, and r = 2m.
    2. Substitute the values into the formula:
    V = (9 × 10^9) × (5 × 10^−6) / 2
    3. Calculate the result:
    V = 22,500 Volts.
    Let me know if you need further clarification!"*
    
    Context: {context}
    Question: {question}
    """

    model = ChatGoogleGenerativeAI(
        model="gemini-2.0-flash",
        temperature=0.3
    )

    prompt = PromptTemplate(
        template=prompt_template,
        input_variables=["context", "question"]
    )
    
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def process_user_question(user_question):
    """Process a user question against the vector store"""
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        # Fix: Added allow_dangerous_deserialization=True
        vector_store = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
        docs = vector_store.similarity_search(user_question, k=4)  # Get top 4 relevant chunks
        
        chain = get_conversational_chain()
        
        # Get the document contents
        doc_contents = [doc.page_content for doc in docs]
        context = "\n\n".join(doc_contents)
        
        response = chain(
            {"input_documents": docs, "context": context, "question": user_question},
            return_only_outputs=True
        )
        
        return response["output_text"]
    except Exception as e:
        return f"Error processing your question: {str(e)}"

# Initialize session state variables
def init_session_state():
    if "files_processed" not in st.session_state:
        st.session_state.files_processed = False
    if "file_info" not in st.session_state:
        st.session_state.file_info = []
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    if "clear_conversation" not in st.session_state:
        st.session_state.clear_conversation = False

def main():
    
    # Initialize session state
    init_session_state()
    
    # Set page configuration
    st.set_page_config(
        page_title="Question Paper Analyzer",
        page_icon="📚",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    # Custom CSS for dark theme
    st.markdown("""
    <style>
    /* Base styling for dark theme */
    .reportview-container {
        background-color: #1E1E1E;
        color: #E0E0E0;
    }
    
    .main .block-container {
        background-color: #1E1E1E;
    }
    
    /* Headers styling */
    .main-header {
        font-size: 2.5rem;
        color: #64FFDA;
        text-align: center;
        margin-bottom: 1rem;
        font-weight: 700;
    }
    
    .sub-header {
        font-size: 1.5rem;
        color: #80CBC4;
        margin-top: 2rem;
        margin-bottom: 1rem;
        font-weight: 600;
    }
    
    /* File box styling */
    .file-box {
        background-color: #2C2C2C;
        border-radius: 10px;
        padding: 20px;
        margin-bottom: 20px;
        border: 1px solid #444444;
    }
    
    /* Chat message styling */
    .chat-message {
        padding: 1.5rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
        display: flex;
    }
    
    .chat-message.user {
        background-color: #263238;
        border-left: 5px solid #64FFDA;
    }
    
    .chat-message.bot {
        background-color: #212121;
        border-left: 5px solid #80CBC4;
    }
    
    .chat-message .avatar {
        width: 20%;
    }
    
    .chat-message .message {
        width: 80%;
    }
    
    /* Question box styling */
    .question-box {
        background-color: #2D3748;
        border-radius: 10px;
        padding: 20px;
        margin: 20px 0;
        border: 1px solid #3C4B5D;
    }
    
    /* Button styling */
    .stButton > button {
        background-color: #64FFDA;
        color: #121212;
        font-weight: bold;
        border: none;
        border-radius: 5px;
        padding: 0.5rem 1rem;
        width: 100%;
    }
    
    .stButton > button:hover {
        background-color: #4DB6AC;
    }
    
    /* File info styling */
    .file-info {
        padding: 10px;
        border-radius: 5px;
        background-color: #2C3E50;
        margin: 5px 0;
        border-left: 3px solid #64FFDA;
    }
    
    /* Upload section styling */
    .upload-section {
        background-color: #2D3748;
        padding: 20px;
        border-radius: 10px;
        border: 1px dashed #4A5568;
    }
    
    /* Progress bar styling */
    .stProgress > div > div > div {
        background-color: #64FFDA;
    }
    
    /* Tabs styling */
    .stTabs [data-baseweb="tab-list"] {
        gap: 24px;
        background-color: #1E1E1E;
    }
    
    .stTabs [data-baseweb="tab"] {
        background-color: #2D3748;
        padding: 10px 20px;
        border-radius: 5px 5px 0 0;
        color: #E0E0E0;
    }
    
    .stTabs [aria-selected="true"] {
        background-color: #64FFDA !important;
        color: #121212 !important;
    }
    
    /* Input fields */
    .stTextInput > div > div > input {
        background-color: #2D3748;
        color: #E0E0E0;
        border: 1px solid #4A5568;
    }
    
    /* Expander styling */
    .streamlit-expanderHeader {
        background-color: #2D3748;
        color: #E0E0E0;
    }
    
    .streamlit-expanderContent {
        background-color: #1E1E1E;
        color: #E0E0E0;
    }
    
    /* Info, warning, error boxes */
    .stAlert {
        background-color: #2D3748;
        color: #E0E0E0;
    }
    
    /* Scrollbar styling */
    ::-webkit-scrollbar {
        width: 8px;
        height: 8px;
    }
    
    ::-webkit-scrollbar-track {
        background: #1E1E1E;
    }
    
    ::-webkit-scrollbar-thumb {
        background-color: #444444;
        border-radius: 10px;
    }
    
    ::-webkit-scrollbar-thumb:hover {
        background-color: #64FFDA;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # App header
    st.markdown('<div class="main-header">📚 Previous Year Question Paper Analyzer</div>', unsafe_allow_html=True)
    
    # Create tabs for different sections
    tab1, tab2, tab3 = st.tabs(["📄 Upload Papers", "❓ Ask Questions", "ℹ️ About"])
    
    with tab1:
        st.markdown('<div class="sub-header">Upload Question Papers</div>', unsafe_allow_html=True)
        
        # Information about supported file formats
        with st.expander("ℹ️ Supported File Formats"):
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("**PDF Files**")
                st.markdown("- Standard PDF documents")
                st.markdown("- Scanned PDFs (with OCR support)")
            with col2:
                st.markdown("**PowerPoint Files**")
                st.markdown("- .pptx (PowerPoint 2007+)")
                st.markdown("- .ppt (Legacy PowerPoint)")
        
        # File upload section
        st.markdown('<div class="upload-section">', unsafe_allow_html=True)
        uploaded_files = st.file_uploader(
            "Drag and drop your question papers here",
            type=["pdf", "pptx", "ppt"],
            accept_multiple_files=True,
            help="Upload question papers in PDF or PowerPoint format"
        )
        
        # Display current files
        if uploaded_files:
            st.markdown("**Selected Files:**")
            for file in uploaded_files:
                st.markdown(f"- 📄 {file.name} ({file.size/1024:.1f} KB)")
            
            # Process button with progress animation
            process_col1, process_col2, process_col3 = st.columns([1, 2, 1])
            with process_col2:
                if st.button("Process Question Papers", key="process_btn", help="Start processing the uploaded files"):
                    # Progress bar for visual feedback
                    progress_bar = st.progress(0)
                    
                    # Process and combine text from all files
                    all_text = ""
                    file_info = []
                    
                    for i, file in enumerate(uploaded_files):
                        # Update progress bar
                        progress = (i / len(uploaded_files)) * 100
                        progress_bar.progress(int(progress))
                        
                        st.markdown(f"Processing **{file.name}**...")
                        file_bytes = file.read()
                        file_extension = file.name.split('.')[-1].lower()
                        
                        # Extract text based on file type
                        if file_extension == 'pdf':
                            text = extract_text_from_pdf(file_bytes)
                        elif file_extension in ['pptx', 'ppt']:
                            text = extract_text_from_ppt(file_bytes)
                        else:
                            st.error(f"Unsupported file type: {file_extension}")
                            continue
                        
                        # Add file metadata to the text
                        file_header = f"\n\n--- DOCUMENT: {file.name} ---\n\n"
                        all_text += file_header + text
                        
                        # Store file info for reference
                        file_info.append({
                            "name": file.name,
                            "type": file_extension,
                            "size": len(file_bytes) // 1024,
                            "text_length": len(text)
                        })
                    
                    # Complete progress bar
                    progress_bar.progress(100)
                    
                    # Create text chunks and vector store
                    if all_text:
                        with st.spinner("Creating text chunks and building vector store..."):
                            text_chunks = get_text_chunks(all_text)
                            st.success(f"✅ Created {len(text_chunks)} text chunks")
                            
                            # Create and save vector store
                            get_vector_store(text_chunks)
                            st.success("✅ Vector store built successfully!")
                        
                        # Save session state to indicate processing is complete
                        st.session_state.files_processed = True
                        st.session_state.file_info = file_info
                        
                        # Success message with animation
                        st.balloons()
                        st.success("🎉 All files processed successfully! You can now ask questions about your question papers.")
                        
                        # Provide guidance on next steps
                        st.info("👉 Head to the 'Ask Questions' tab to start analyzing your question papers.")
                    else:
                        st.error("❌ No text could be extracted from the uploaded files.")
        else:
            st.info("Please upload at least one question paper file to begin.")
        st.markdown('</div>', unsafe_allow_html=True)
        
        # Display file information if available
        if st.session_state.files_processed and st.session_state.file_info:
            st.markdown('<div class="sub-header">Processed Files</div>', unsafe_allow_html=True)
            for file in st.session_state.file_info:
                st.markdown(f"""
                <div class="file-info">
                    <b>📄 {file['name']}</b><br>
                    Format: {file['type'].upper()} | Size: {file['size']} KB | Text Length: {file.get('text_length', 'N/A')} characters
                </div>
                """, unsafe_allow_html=True)
    
    with tab2:
        st.markdown('<div class="sub-header">Ask Questions About Your Papers</div>', unsafe_allow_html=True)
        
        if not st.session_state.files_processed:
            st.warning("⚠️ Please upload and process files in the 'Upload Papers' tab before asking questions.")
        else:
            # Example questions to help users get started
            with st.expander("📋 Example Questions You Can Ask"):
                st.markdown("""
                - "What topics appeared most frequently across all papers?"
                - "Show me all questions related to [specific topic] from the 2023 paper."
                - "What is the answer to Question 5(b) in the 2022 paper?"
                - "Compare the difficulty level of questions across different years."
                - "What chapters should I focus on based on previous year trends?"
                """)
            
            # Check if we need to clear the conversation before rendering the input widget
            if st.session_state.clear_conversation:
                # Reset the flag
                st.session_state.clear_conversation = False
                # Clear the key from session state completely
                if "question_input" in st.session_state:
                    del st.session_state["question_input"]
            
            # Question input with chat-like interface
            st.markdown('<div class="question-box">', unsafe_allow_html=True)
            user_question = st.text_input(
                "Ask a question about your question papers:",
                key="question_input",
                placeholder="Type your question here and press Enter or click 'Ask'",
                help="Ask any question about the papers you've uploaded"
            )
            
            # Create a clean layout for the ask button
            col1, col2, col3 = st.columns([3, 1, 3])
            with col2:
                ask_button = st.button("Ask", key="ask_button")
            st.markdown('</div>', unsafe_allow_html=True)
            
            # Process the question when the button is clicked or Enter is pressed
            if (ask_button or user_question != "") and user_question:
                # Add user question to chat history
                st.session_state.chat_history.append({"role": "user", "content": user_question})
                
                # Process the question
                with st.spinner("Analyzing question papers to find your answer..."):
                    response = process_user_question(user_question)
                    # Add bot response to chat history
                    st.session_state.chat_history.append({"role": "bot", "content": response})
            
            # Display chat history
            st.markdown('<div class="sub-header">Conversation History</div>', unsafe_allow_html=True)
            
            if len(st.session_state.chat_history) == 0:
                st.info("Your conversation will appear here. Ask a question to get started!")
            else:
                # Track the last displayed user message to avoid duplicates
                last_user_message = None
                
                for message in st.session_state.chat_history:
                    if message["role"] == "user":
                        # Skip if this message is identical to the previous user message
                        if message["content"] == last_user_message:
                            continue
                        
                        # Display the message and update last_user_message
                        st.markdown(f"""
                        <div class="chat-message user">
                            <div class="avatar">
                                <strong>You:</strong>
                            </div>
                            <div class="message">
                                {message["content"]}
                            </div>
                        </div>
                        """, unsafe_allow_html=True)
                        
                        last_user_message = message["content"]
                    else:
                        st.markdown(f"""
                        <div class="chat-message bot">
                            <div class="avatar">
                                <strong>AI:</strong>
                            </div>
                            <div class="message">
                                {message["content"]}
                            </div>
                        </div>
                        """, unsafe_allow_html=True)
            
            # Clear conversation button
            if st.session_state.chat_history:
                col1, col2, col3 = st.columns([3, 1, 3])
                with col2:
                    if st.button("Clear Conversation", key="clear_btn"):
                        # Clear chat history
                        st.session_state.chat_history = []
                        # Set flag to clear input on next rerun
                        st.session_state.clear_conversation = True
                        st.rerun()

    
    with tab3:
        st.markdown('<div class="sub-header">About This Tool</div>', unsafe_allow_html=True)

        # Tool description
        st.markdown("""
        The **Previous Year Question Paper Analyzer** is an advanced AI-powered tool designed to help students, educators, and researchers analyze patterns and trends in question papers. This tool enables users to:
        
        - **Upload multiple question papers** in PDF or PowerPoint format
        - **Extract and analyze text** from these papers, including scanned documents using OCR
        - **Ask questions** about the content, patterns, and trends in these papers
        - **Get detailed answers** with proper citations to the source documents
        
        This tool is particularly useful for:
        - **Students** preparing for exams by understanding question patterns
        - **Teachers** analyzing exam patterns for more effective teaching
        - **Researchers** studying educational assessment trends
        - **Educational institutions** evaluating examination quality and consistency
        """)
        
        # Features section
        st.markdown('<div class="sub-header">Key Features</div>', unsafe_allow_html=True)
        
        col1, col2 = st.columns(2)
        with col1:
            st.markdown("""
            - **Multi-document analysis**: Process several years of papers at once
            - **Intelligent text extraction**: OCR capabilities for scanned documents
            - **Natural language queries**: Ask questions in plain English
            - **Pattern recognition**: Identify recurring topics and question formats
            """)
        
        with col2:
            st.markdown("""
            - **Citation support**: Answers reference specific papers and questions
            - **Comprehensive insights**: Analyze difficulty trends and mark distributions
            - **User-friendly interface**: Simple upload and query process
            - **Educational focus**: Designed specifically for academic analysis
            """)
        
        # How to use section
        st.markdown('<div class="sub-header">How to Use</div>', unsafe_allow_html=True)
        
        st.markdown("""
        1. **Upload Papers**: Go to the 'Upload Papers' tab and upload your question papers (PDF or PPT format)
        2. **Process Files**: Click the 'Process Question Papers' button to extract and analyze the content
        3. **Ask Questions**: Navigate to the 'Ask Questions' tab and type your query about the papers
        4. **Review Answers**: Get AI-generated answers based on the content of your uploaded papers
        """)
        
        # Requirements note
        with st.expander("System Requirements"):
            st.markdown("""
            - **Google API Key**: Required for AI processing (set as environment variable)
            - **PDF OCR Support** (optional): For better processing of scanned documents
              - pdf2image
              - pytesseract
              - poppler
            - **Internet Connection**: Required for AI processing
            """)

# Run the main application
if __name__ == "__main__":
    main()